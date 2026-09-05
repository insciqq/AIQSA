"""Synthetic inputs and independent artifact checks for the opt-in paid smoke.

Run only in the networkless workspace-guest image, never execute model code.
Fixtures are ordinary, already-filled Office documents, without template tags.
"""
import base64
import io
import json
import re
import sys
import uuid
import zipfile
from decimal import Decimal

from docx import Document
from docx.shared import Pt
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor


def encoded(document):
    buffer = io.BytesIO()
    document.save(buffer)
    return base64.b64encode(buffer.getvalue()).decode("ascii")


def presentation(title, lines):
    result = Presentation()
    for heading, body in [(title, "Northwind Studio"), *lines]:
        slide = result.slides.add_slide(result.slide_layouts[1])
        slide.shapes.title.text = heading
        slide.placeholders[1].text = body
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = RGBColor.from_string("1E40AF")
    return result


def fixtures():
    leave = Document()
    leave.styles["Normal"].font.name = "Arial"
    leave.styles["Normal"].font.size = Pt(11)
    header = leave.add_table(rows=2, cols=1)
    header.cell(0, 0).text = "Генеральному директору ООО Север"
    header.cell(1, 0).text = "от Иванова Ивана Ивановича, аналитика"
    leave.add_heading("Заявление на отпуск", level=1)
    paragraph = leave.add_paragraph("Прошу предоставить мне ежегодный оплачиваемый отпуск с ")
    paragraph.add_run("1 июня 2025 года").underline = True
    paragraph.add_run(" по ")
    paragraph.add_run("14 июня 2025 года").underline = True
    paragraph.add_run(" включительно.")
    leave.add_paragraph("Дата заявления: 20 мая 2025 года")
    leave.add_paragraph("Подпись: ____________________")
    leave.sections[0].footer.paragraphs[0].text = "Кадровая служба · ООО Север"

    report = Document()
    report.add_heading("Northwind Studio", 0)
    report.add_heading("Monthly management report — August 2026", 1)
    report.add_paragraph("Revenue was 12000, costs were 8000 and profit was 4000.")
    table = report.add_table(rows=1, cols=2)
    table.style = "Light Shading Accent 1"
    table.rows[0].cells[0].text, table.rows[0].cells[1].text = "Metric", "Amount"
    for label, amount in [("Revenue", "12000"), ("Costs", "8000"), ("Profit", "4000")]:
        row = table.add_row().cells
        row[0].text, row[1].text = label, amount
    report.add_heading("Management actions", 1)
    report.add_paragraph("Review departmental costs at the next monthly meeting.")

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    for row in [["Northwind Studio"], ["August 2026"], [], ["Metric", "Amount"], ["Revenue", 12000], ["Costs", 8000], ["Profit", 4000]]:
        sheet.append(row)
    for cell in sheet[4]:
        cell.fill = PatternFill("solid", fgColor="1E40AF")
        cell.font = Font(color="FFFFFF", bold=True)
    sheet.column_dimensions["A"].width = 28
    sheet.column_dimensions["B"].width = 18
    slides = presentation("Monthly review — August 2026", [
        ("Financial results", "Revenue 12000\nCosts 8000\nProfit 4000"),
        ("Management actions", "Review departmental costs next month.")
    ])
    stock = Workbook()
    for row in [["sku", "product", "quantity", "unit_price", "line_total"], ["0007", "Notebook", 12, 4.5, 54], ["0042", "Pen", 30, 1.2, 36], ["0105", "Folder", 8, 2.75, 22]]:
        stock.active.append(row)
    return {"leave-application.docx": encoded(leave), "report-example.docx": encoded(report),
            "report-example.xlsx": encoded(workbook), "meeting-example.pptx": encoded(slides), "stock-source.xlsx": encoded(stock)}


def package(raw):
    value = zipfile.ZipFile(io.BytesIO(raw))
    assert value.testzip() is None
    assert "[Content_Types].xml" in value.namelist()
    return value


def words(document):
    text = [paragraph.text for paragraph in document.paragraphs]
    text += [cell.text for table in document.tables for row in table.rows for cell in row.cells]
    text += [paragraph.text for section in document.sections for paragraph in section.footer.paragraphs]
    return "\n".join(text)


def slides_text(document):
    values = []
    for slide in document.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                values.append(shape.text)
            if shape.has_table:
                values.extend(cell.text for row in shape.table.rows for cell in row.cells)
            if shape.has_chart:
                for plot in shape.chart.plots:
                    values.extend(str(category.label) for category in plot.categories)
                for series in shape.chart.series:
                    values.extend(str(value) for value in series.values)
    return "\n".join(values)


def numeric_values(text):
    # Office authors may use grouping spaces and a decimal comma in prose.
    normalized = re.sub(r"(?<=\d)[,\u00a0\u202f ](?=\d{3}(?:\D|$))", "", text)
    return {Decimal(match.replace(",", ".")) for match in re.findall(r"\d+(?:[.,]\d+)?", normalized)}


def check(case, files):
    opened = {}
    for name, raw in files.items():
        if name.endswith((".xlsx", ".docx", ".pptx")):
            package(raw).close()
        if name.endswith(".xlsx"):
            opened[name] = load_workbook(io.BytesIO(raw))
        elif name.endswith(".docx"):
            opened[name] = Document(io.BytesIO(raw))
        elif name.endswith(".pptx"):
            opened[name] = Presentation(io.BytesIO(raw))

    if case in ("office_stock_blue", "office_stock_purple", "office_stock_notes"):
        suffix = case.removeprefix("office_stock_")
        book = opened[f"stock-{suffix}.xlsx"]
        sheet = book.active
        assert [sheet.cell(row, 1).value for row in range(2, 5)] == ["0007", "0042", "0105"]
        assert [sheet.cell(row, 2).value for row in range(2, 5)] == ["Notebook", "Pen", "Folder"]
        assert [sheet.cell(row, 3).value for row in range(2, 5)] == [12, 30, 8]
        assert [Decimal(str(sheet.cell(row, 4).value)) for row in range(2, 5)] == [Decimal("4.50"), Decimal("1.20"), Decimal("2.75")]
        for row, amount in [(2, "54"), (3, "36"), (4, "22")]:
            value = sheet.cell(row, 5).value
            assert (isinstance(value, str) and value.replace(" ", "").upper() == f"=C{row}*D{row}") or Decimal(str(value)) == Decimal(amount)
        color = "7C3AED" if suffix == "purple" else "1E40AF"
        assert sheet["A1"].fill.fgColor.rgb[-6:].upper() == color
        assert sheet.freeze_panes and sheet.auto_filter.ref
        assert uuid.UUID(book.properties.identifier).version == 4
        if suffix == "notes":
            assert sheet["F1"].value == "Notes" and sheet["F2"].value == "Reorder in October"
        if "previous.xlsx" in opened:
            previous = opened["previous.xlsx"]
            assert previous.properties.identifier == book.properties.identifier
            if suffix == "purple":
                assert list(previous.active.values) == list(sheet.values)
    elif case == "office_slides":
        slides = opened["stock-review.pptx"]
        assert len(slides.slides) == 4
        text = slides_text(slides)
        assert all(name in text for name in ("Notebook", "Pen", "Folder"))
        assert Decimal("112") in numeric_values(text)
    elif case == "office_slide_edit":
        before, after = opened["stock-review.pptx"], opened["stock-review-revised.pptx"]
        assert len(before.slides) == len(after.slides) == 4
        third = "\n".join(shape.text for shape in after.slides[2].shapes if shape.has_text_frame)
        assert "Updated purchasing plan" in third and "Reorder 20 notebooks" in third
        # Compare the exported XML bytes. python-pptx's XmlString overloads ==
        # with a line parser that can reject valid multiline slide text.
        with zipfile.ZipFile(io.BytesIO(files["stock-review.pptx"])) as original, \
                zipfile.ZipFile(io.BytesIO(files["stock-review-revised.pptx"])) as revised:
            for index in [1, 2, 4]:
                path = f"ppt/slides/slide{index}.xml"
                assert original.read(path) == revised.read(path)
    elif case in ("office_docx", "office_leave"):
        output = opened["leave-updated.docx" if case == "office_docx" else "leave-september.docx"]
        source = opened["leave-application.docx"]
        text = words(output)
        assert len(output.tables) == len(source.tables) == 1
        assert len(output.tables[0].rows) == 2 and len(output.sections) == len(source.sections)
        assert "ООО Север" in text and "Кадровая служба" in text
        assert "Иванов" in text and "2025" not in text
        for expected in (["5 августа 2026", "18 августа 2026", "25 июля 2026"] if case == "office_docx" else ["1 сентября 2027", "15 сентября 2027", "20 августа 2027"]):
            assert expected in text
        assert any(run.underline for paragraph in output.paragraphs for run in paragraph.runs)
        assert "Подпись" in text and "____" in text
    elif case in ("office_month_september", "office_month_october"):
        period = "September 2026" if case.endswith("september") else "October 2026"
        name = "september" if case.endswith("september") else "october"
        totals = {Decimal(value) for value in ([20000, 13000, 7000] if name == "september" else [26000, 16000, 10000])}
        book = opened[f"{name}-report.xlsx"]
        report = opened[f"{name}-report.docx"]
        slides = opened[f"{name}-meeting.pptx"]
        grid = "\n".join(str(cell.value) for sheet in book for row in sheet for cell in row if cell.value is not None)
        for text in [grid, words(report), slides_text(slides)]:
            assert "Northwind Studio" in text and period in text
            assert totals.issubset(numeric_values(text))
            assert "August 2026" not in text
            if name == "october":
                assert "September 2026" not in text
        assert report.tables and len(slides.slides) == 3
    else:
        raise AssertionError("unknown_office_case")


data = json.load(sys.stdin)
if data["case"] == "office_fixtures":
    print(json.dumps({"files": fixtures()}))
else:
    check(data["case"], {key: base64.b64decode(value) for key, value in data["files"].items()})
    print(json.dumps({"oraclePassed": True}))
